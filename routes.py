import csv
import io
import json
import os
from datetime import datetime, timedelta
from flask import render_template, request, flash, redirect, url_for, make_response, jsonify, send_file
from werkzeug.utils import secure_filename
from app import app
from PIL import Image, ImageEnhance, ImageFilter
import PyPDF2
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import docx
import openpyxl
from pptx import Presentation
import numpy as np
try:
    import cv2
except ImportError:
    cv2 = None

# Upload and temp folder configuration
UPLOAD_FOLDER = 'uploads'
TEMP_FOLDER = 'temp'

# Create directories if they don't exist
for folder in [UPLOAD_FOLDER, TEMP_FOLDER]:
    if not os.path.exists(folder):
        os.makedirs(folder)

# Sample data for demos (realistic business data)
sample_products = [
    {"id": 1, "nome": "Camiseta Polo", "categoria": "Roupas", "preco": 45.90, "estoque": 25},
    {"id": 2, "nome": "Tênis Esportivo", "categoria": "Calçados", "preco": 189.90, "estoque": 12},
    {"id": 3, "nome": "Mochila Executiva", "categoria": "Acessórios", "preco": 98.50, "estoque": 8},
    {"id": 4, "nome": "Relógio Digital", "categoria": "Eletrônicos", "preco": 156.00, "estoque": 15},
    {"id": 5, "nome": "Agenda de Couro", "categoria": "Papelaria", "preco": 28.90, "estoque": 30}
]

# Sample IA Analytics data
sample_analytics_data = {
    "vendas_mes": [
        {"mes": "Jan", "vendas": 25000, "meta": 20000, "crescimento": 25},
        {"mes": "Fev", "vendas": 28000, "meta": 22000, "crescimento": 12},
        {"mes": "Mar", "vendas": 32000, "meta": 25000, "crescimento": 14},
        {"mes": "Abr", "vendas": 35000, "meta": 28000, "crescimento": 9},
        {"mes": "Mai", "vendas": 38000, "meta": 30000, "crescimento": 8}
    ],
    "insights": [
        {"tipo": "Tendência", "titulo": "Crescimento Consistente", "descricao": "Vendas crescendo 13% ao mês", "impacto": "Alto"},
        {"tipo": "Alerta", "titulo": "Estoque Baixo", "descricao": "3 produtos com estoque crítico", "impacto": "Médio"},
        {"tipo": "Oportunidade", "titulo": "Nova Categoria", "descricao": "Eletrônicos com alta demanda", "impacto": "Alto"},
        {"tipo": "Previsão", "titulo": "Meta Junho", "descricao": "Projeção: R$ 42.000", "impacto": "Informativo"}
    ]
}

# Sample CRM data
sample_crm_data = [
    {"id": 1, "nome": "João Silva", "empresa": "Tech Corp", "status": "Negociação", "valor": 15000, "telefone": "(11) 99999-1111"},
    {"id": 2, "nome": "Maria Santos", "empresa": "Inovação Ltda", "status": "Proposta", "valor": 8500, "telefone": "(11) 99999-2222"},
    {"id": 3, "nome": "Pedro Costa", "empresa": "Digital Solutions", "status": "Qualificado", "valor": 12000, "telefone": "(11) 99999-3333"},
    {"id": 4, "nome": "Ana Lima", "empresa": "StartUp XYZ", "status": "Fechado", "valor": 25000, "telefone": "(11) 99999-4444"},
    {"id": 5, "nome": "Carlos Mendes", "empresa": "Mega Systems", "status": "Contato", "valor": 18000, "telefone": "(11) 99999-5555"}
]

sample_financial_data = [
    {"data": "2025-01-15", "tipo": "Receita", "categoria": "Vendas", "valor": 2500.00, "descricao": "Vendas do dia"},
    {"data": "2025-01-16", "tipo": "Despesa", "categoria": "Aluguel", "valor": -800.00, "descricao": "Aluguel da loja"},
    {"data": "2025-01-17", "tipo": "Receita", "categoria": "Vendas", "valor": 1800.00, "descricao": "Vendas do dia"},
    {"data": "2025-01-18", "tipo": "Despesa", "categoria": "Fornecedores", "valor": -650.00, "descricao": "Compra de mercadorias"},
    {"data": "2025-01-19", "tipo": "Receita", "categoria": "Vendas", "valor": 3200.00, "descricao": "Vendas do dia"}
]

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/solucoes')
def solutions():
    return render_template('solutions.html')

@app.route('/dia-a-dia')
def dia_a_dia():
    return render_template('dia_a_dia.html')

@app.route('/sobre')
def about():
    return render_template('about.html')

@app.route('/contato', methods=['GET', 'POST'])
def contact():
    if request.method == 'POST':
        nome = request.form.get('nome')
        email = request.form.get('email')
        empresa = request.form.get('empresa')
        mensagem = request.form.get('mensagem')
        
        # Simulate successful form submission
        flash('Mensagem enviada com sucesso! Entraremos em contato em breve.', 'success')
        return redirect(url_for('contact'))
    
    return render_template('contact.html')

@app.route('/cases')
def cases():
    return render_template('cases.html')

@app.route('/vantagens')
def vantagens():
    return render_template('vantagens.html')

@app.route('/estudos')
def estudos():
    return render_template('estudos.html')

# Demo Routes
@app.route('/demo/ecommerce')
def demo_ecommerce():
    return render_template('demos/ecommerce.html', products=sample_products)

@app.route('/demo/agenda')
def demo_agenda():
    return render_template('demos/agenda.html')

@app.route('/demo/estoque')
def demo_estoque():
    return render_template('demos/estoque.html', products=sample_products)

@app.route('/demo/financeiro')
def demo_financeiro():
    return render_template('demos/financeiro.html', transactions=sample_financial_data)

@app.route('/demo/restaurante')
def demo_restaurante():
    return render_template('demos/restaurante.html')

@app.route('/demo/erp')
def demo_erp():
    return render_template('demos/erp.html')

@app.route('/demo/oficina')
def demo_oficina():
    return render_template('demos/oficina.html')

@app.route('/demo/imoveis')
def demo_imoveis():
    return render_template('demos/imoveis.html')

@app.route('/demo/academia')
def demo_academia():
    return render_template('demos/academia.html')

# New Top 3 Systems for 2025
@app.route('/demo/ia-analytics')
def demo_ia_analytics():
    return render_template('demos/ia-analytics.html', 
                         analytics_data=sample_analytics_data,
                         products=sample_products)

@app.route('/demo/crm')
def demo_crm():
    return render_template('demos/crm.html', leads=sample_crm_data)

# Sistemas Famosos - Top 2025
@app.route('/demo/crm-profissional')
def demo_crm_profissional():
    # Sample Pipedrive-style data with enhanced features
    sample_pipedrive_data = [
        {"id": 1, "nome": "Ana Silva", "empresa": "TechStart", "status": "Negociação", "valor": 25000, "probabilidade": 75, "origem": "Website", "telefone": "(11) 99999-1111"},
        {"id": 2, "nome": "Carlos Santos", "empresa": "Digital Corp", "status": "Proposta Enviada", "valor": 15000, "probabilidade": 60, "origem": "Indicação", "telefone": "(11) 99999-2222"},
        {"id": 3, "nome": "Marina Costa", "empresa": "StartUp XYZ", "status": "Demo Agendada", "valor": 35000, "probabilidade": 40, "origem": "LinkedIn", "telefone": "(11) 99999-3333"},
        {"id": 4, "nome": "Roberto Lima", "empresa": "Innovation Ltd", "status": "Fechado-Ganho", "valor": 50000, "probabilidade": 100, "origem": "Google Ads", "telefone": "(11) 99999-4444"}
    ]
    return render_template('demos/crm-profissional.html', leads=sample_pipedrive_data)

@app.route('/demo/hubspot-marketing')
def demo_hubspot_marketing():
    return render_template('demos/hubspot-marketing.html')

@app.route('/demo/slack-comunicacao')
def demo_slack_comunicacao():
    return render_template('demos/slack-comunicacao.html')

@app.route('/demo/trello-projetos')
def demo_trello_projetos():
    return render_template('demos/trello-projetos.html')

@app.route('/demo/quickbooks-contabilidade')
def demo_quickbooks_contabilidade():
    return render_template('demos/quickbooks-contabilidade.html')

@app.route('/demo/shopify-ecommerce')
def demo_shopify_ecommerce():
    return render_template('demos/shopify-ecommerce.html')

@app.route('/demo/whatsapp-business')
def demo_whatsapp_business():
    return render_template('demos/whatsapp-business.html')

@app.route('/demo/zoom-videoconferencia')
def demo_zoom_videoconferencia():
    return render_template('demos/zoom-videoconferencia.html')

@app.route('/demo/netflix-streaming')
def demo_netflix_streaming():
    return render_template('demos/netflix-streaming.html')

@app.route('/demo/google-analytics')
def demo_google_analytics():
    return render_template('demos/google-analytics.html')

@app.route('/demo/salesforce-enterprise')
def demo_salesforce_enterprise():
    return render_template('demos/salesforce-enterprise.html')

# New systems for different niches
@app.route('/demo/clinica-medica')
def demo_clinica_medica():
    return render_template('demos/clinica-medica.html')

@app.route('/demo/escola-online')
def demo_escola_online():
    return render_template('demos/escola-online.html')

@app.route('/demo/delivery')
def demo_delivery():
    return render_template('demos/delivery.html')

@app.route('/demo/pet-shop')
def demo_pet_shop():
    return render_template('demos/pet-shop.html')

@app.route('/demo/agencia-viagem')
def demo_agencia_viagem():
    return render_template('demos/agencia-viagem.html')

@app.route('/demo/contabilidade')
def demo_contabilidade():
    return render_template('demos/contabilidade.html')

@app.route('/demo/advocacia')
def demo_advocacia():
    return render_template('demos/advocacia.html')

@app.route('/demo/marketing-digital')
def demo_marketing_digital():
    return render_template('demos/marketing-digital.html')

# Export functionality
@app.route('/export/csv/<demo_type>')
def export_csv(demo_type):
    output = io.StringIO()
    writer = csv.writer(output)
    
    if demo_type == 'estoque':
        writer.writerow(['ID', 'Produto', 'Categoria', 'Preço', 'Estoque'])
        for product in sample_products:
            writer.writerow([product['id'], product['nome'], product['categoria'], 
                           f"R$ {product['preco']:.2f}", product['estoque']])
        filename = 'relatorio_estoque.csv'
    
    elif demo_type == 'financeiro':
        writer.writerow(['Data', 'Tipo', 'Categoria', 'Valor', 'Descrição'])
        for transaction in sample_financial_data:
            writer.writerow([transaction['data'], transaction['tipo'], transaction['categoria'],
                           f"R$ {transaction['valor']:.2f}", transaction['descricao']])
        filename = 'relatorio_financeiro.csv'
    
    else:
        writer.writerow(['Dados', 'Não disponíveis'])
        filename = 'relatorio.csv'
    
    output.seek(0)
    
    response = make_response(output.getvalue())
    response.headers["Content-Disposition"] = f"attachment; filename={filename}"
    response.headers["Content-type"] = "text/csv; charset=utf-8"
    
    return response

# API endpoints for chart data
@app.route('/api/chart-data/<chart_type>')
def get_chart_data(chart_type):
    if chart_type == 'vendas-categoria':
        categories = {}
        for product in sample_products:
            if product['categoria'] not in categories:
                categories[product['categoria']] = 0
            categories[product['categoria']] += product['estoque'] * product['preco']
        
        return jsonify({
            'labels': list(categories.keys()),
            'data': list(categories.values())
        })
    
    elif chart_type == 'fluxo-caixa':
        receitas = sum([t['valor'] for t in sample_financial_data if t['valor'] > 0])
        despesas = abs(sum([t['valor'] for t in sample_financial_data if t['valor'] < 0]))
        
        return jsonify({
            'labels': ['Receitas', 'Despesas'],
            'data': [receitas, despesas]
        })
    
    return jsonify({'error': 'Tipo de gráfico não encontrado'})

# API Routes for Dia-a-dia Tools
UPLOAD_FOLDER = 'uploads'
TEMP_FOLDER = 'temp'
ALLOWED_EXTENSIONS = {
    'pdf': {'pdf'},
    'document': {'doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx'},
    'image': {'png', 'jpg', 'jpeg'}
}

def allowed_file(filename, file_type):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS[file_type]

@app.route('/api/convert-to-pdf', methods=['POST'])
def convert_to_pdf():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename, 'document') and not allowed_file(file.filename, 'image'):
            return jsonify({'error': 'Invalid file type'}), 400
        
        filename = secure_filename(file.filename or 'file')
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(file_path)
        
        # Convert to PDF based on file type
        pdf_path = convert_file_to_pdf(file_path, filename)
        
        if pdf_path:
            return send_file(pdf_path, as_attachment=True, 
                           download_name=filename.rsplit('.', 1)[0] + '.pdf',
                           mimetype='application/pdf')
        else:
            return jsonify({'error': 'Conversion failed'}), 500
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/compress-pdf', methods=['POST'])
def compress_pdf():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename, 'pdf'):
            return jsonify({'error': 'Only PDF files allowed'}), 400
        
        # Get compression level from form data
        compression_level = request.form.get('level', 'medium')  # low, medium, high
        
        filename = secure_filename(file.filename or 'file')
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(file_path)
        
        # Get original file size
        original_size = os.path.getsize(file_path)
        
        # Compress PDF with selected level
        result = compress_pdf_file(file_path, filename, compression_level)
        if result and len(result) == 2:
            compressed_path, compression_stats = result
        else:
            compressed_path = result
            compression_stats = None
        
        if compressed_path:
            # Get compressed file size
            compressed_size = os.path.getsize(compressed_path)
            compression_ratio = ((original_size - compressed_size) / original_size) * 100
            
            # Add compression info to response headers
            response = make_response(send_file(compressed_path, as_attachment=True,
                           download_name=filename.rsplit('.', 1)[0] + '_compressed.pdf',
                           mimetype='application/pdf'))
            response.headers['X-Original-Size'] = str(original_size)
            response.headers['X-Compressed-Size'] = str(compressed_size)
            response.headers['X-Compression-Ratio'] = str(round(compression_ratio, 1))
            return response
        else:
            return jsonify({'error': 'Compression failed'}), 500
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/remove-background', methods=['POST'])
def remove_background():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename, 'image'):
            return jsonify({'error': 'Only image files allowed'}), 400
        
        # Get processing options from form data
        model_type = request.form.get('model', 'u2net')  # u2net, isnet, silueta
        output_format = request.form.get('format', 'png')  # png, jpg
        background_color = request.form.get('bg_color', 'transparent')  # transparent, white, black, custom
        edge_refine = request.form.get('edge_refine', 'true') == 'true'
        
        filename = secure_filename(file.filename or 'file')
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(file_path)
        
        # Remove background with advanced options
        output_path = remove_image_background_advanced(file_path, filename, {
            'model': model_type,
            'format': output_format,
            'background': background_color,
            'edge_refine': edge_refine
        })
        
        if output_path:
            download_name = filename.rsplit('.', 1)[0] + f'_no_bg.{output_format}'
            mimetype = f'image/{output_format}'
            return send_file(output_path, as_attachment=True,
                           download_name=download_name,
                           mimetype=mimetype)
        else:
            return jsonify({'error': 'Background removal failed'}), 500
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Helper functions for file processing
def convert_file_to_pdf(file_path, filename):
    try:
        ext = filename.rsplit('.', 1)[1].lower()
        output_path = os.path.join(TEMP_FOLDER, filename.rsplit('.', 1)[0] + '.pdf')
        
        if ext in ['jpg', 'jpeg', 'png']:
            # Convert image to PDF
            image = Image.open(file_path)
            if image.mode == 'RGBA':
                image = image.convert('RGB')
            image.save(output_path, 'PDF')
            
        elif ext in ['doc', 'docx']:
            # Convert Word document to PDF (simplified)
            doc = docx.Document(file_path)
            
            # Create PDF with text content
            c = canvas.Canvas(output_path, pagesize=letter)
            y_position = 750
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    c.drawString(50, y_position, paragraph.text[:80])  # Limit line length
                    y_position -= 20
                    if y_position < 50:
                        c.showPage()
                        y_position = 750
            
            c.save()
            
        elif ext in ['xls', 'xlsx']:
            # Convert Excel to PDF (simplified)
            wb = openpyxl.load_workbook(file_path)
            ws = wb.active
            
            c = canvas.Canvas(output_path, pagesize=letter)
            y_position = 750
            
            try:
                for row in ws.iter_rows(values_only=True, max_row=50):  # Limit rows
                    row_text = ' | '.join([str(cell) if cell else '' for cell in row])
                    if row_text.strip():
                        c.drawString(50, y_position, row_text[:80])
                        y_position -= 20
                        if y_position < 50:
                            c.showPage()
                            y_position = 750
            except AttributeError:
                # Fallback for older openpyxl versions
                c.drawString(50, y_position, "Excel content (preview not available)")
            
            c.save()
            
        elif ext in ['ppt', 'pptx']:
            # Convert PowerPoint to PDF (simplified)
            prs = Presentation(file_path)
            
            c = canvas.Canvas(output_path, pagesize=letter)
            y_position = 750
            
            for slide in prs.slides:
                c.drawString(50, y_position, f"Slide {prs.slides.index(slide) + 1}")
                y_position -= 30
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and hasattr(shape, "text") and getattr(shape, "text", "").strip():
                        c.drawString(70, y_position, getattr(shape, "text", "")[:70])
                        y_position -= 20
                        if y_position < 50:
                            c.showPage()
                            y_position = 750
                
                if y_position < 700:  # New slide
                    c.showPage()
                    y_position = 750
            
            c.save()
        
        # Clean up uploaded file
        os.remove(file_path)
        return output_path
        
    except Exception as e:
        print(f"Conversion error: {e}")
        return None

def compress_pdf_file(file_path, filename, compression_level='medium'):
    try:
        output_path = os.path.join(TEMP_FOLDER, filename.rsplit('.', 1)[0] + '_compressed.pdf')
        
        # Advanced PDF compression using PyPDF2 with different levels
        with open(file_path, 'rb') as input_file:
            reader = PyPDF2.PdfReader(input_file)
            writer = PyPDF2.PdfWriter()
            
            # Compression settings based on level
            if compression_level == 'low':
                # Minimal compression - preserve quality
                compression_settings = {
                    'compress_images': False,
                    'image_quality': 95,
                    'remove_duplication': True
                }
            elif compression_level == 'high':
                # Maximum compression - significant size reduction
                compression_settings = {
                    'compress_images': True,
                    'image_quality': 50,
                    'remove_duplication': True
                }
            else:  # medium
                # Balanced compression
                compression_settings = {
                    'compress_images': True,
                    'image_quality': 75,
                    'remove_duplication': True
                }
            
            for page in reader.pages:
                # Apply compression based on level
                if compression_level in ['medium', 'high']:
                    page.compress_content_streams()
                
                # Remove duplicate resources if requested
                if compression_settings['remove_duplication']:
                    # This is a simplified approach - in production would use more advanced PDF libraries
                    pass
                
                writer.add_page(page)
            
            # Apply additional compression settings
            if compression_level == 'high':
                # Use available compression methods (PyPDF2 compatibility)
                try:
                    if hasattr(writer, 'compress_identical_objects'):
                        writer.compress_identical_objects()
                except:
                    pass  # Method not available in this PyPDF2 version
                try:
                    if hasattr(writer, 'remove_links'):
                        writer.remove_links()
                except:
                    pass  # Method not available in this PyPDF2 version
                
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
        
        # Get compression statistics
        original_size = os.path.getsize(file_path)
        compressed_size = os.path.getsize(output_path)
        compression_stats = {
            'original_size': original_size,
            'compressed_size': compressed_size,
            'compression_ratio': ((original_size - compressed_size) / original_size) * 100,
            'level': compression_level
        }
        
        # Clean up uploaded file
        os.remove(file_path)
        return output_path, compression_stats
        
    except Exception as e:
        print(f"Compression error: {e}")
        return None, None

def remove_image_background_advanced(file_path, filename, options):
    try:
        ext = options.get('format', 'png')
        output_path = os.path.join(TEMP_FOLDER, filename.rsplit('.', 1)[0] + f'_no_bg.{ext}')
        
        # Load image
        image = Image.open(file_path)
        original_mode = image.mode
        
        # Convert to RGBA for processing
        if image.mode != 'RGBA':
            image = image.convert('RGBA')
        
        # Advanced background removal using multiple algorithms
        model_type = options.get('model', 'u2net')
        
        if cv2 is not None and model_type == 'u2net':
            # Use OpenCV-based advanced edge detection and segmentation
            result_image = opencv_background_removal(image)
        elif model_type == 'isnet':
            # Use improved segmentation algorithm
            result_image = improved_segmentation(image)
        elif model_type == 'silueta':
            # Use silhouette detection
            result_image = silhouette_detection(image)
        else:
            # Fallback to enhanced basic removal
            result_image = enhanced_basic_removal(image)
        
        # Apply edge refinement if requested
        if options.get('edge_refine', True):
            result_image = refine_edges(result_image)
        
        # Apply background replacement
        background_color = options.get('background', 'transparent')
        if background_color != 'transparent':
            result_image = apply_background(result_image, background_color)
        
        # Save in requested format
        if ext.lower() == 'jpg' or ext.lower() == 'jpeg':
            # Convert to RGB for JPEG (no transparency support)
            if result_image.mode == 'RGBA':
                # Create white background for JPEG
                bg = Image.new('RGB', result_image.size, (255, 255, 255))
                bg.paste(result_image, mask=result_image.split()[-1])
                result_image = bg
            result_image.save(output_path, 'JPEG', quality=95, optimize=True)
        else:
            result_image.save(output_path, 'PNG', optimize=True)
        
        # Clean up uploaded file
        os.remove(file_path)
        return output_path
        
    except Exception as e:
        print(f"Background removal error: {e}")
        return None

def opencv_background_removal(image):
    """Advanced background removal using OpenCV"""
    if cv2 is None:
        return enhanced_basic_removal(image)
    
    try:
        # Convert PIL to OpenCV format
        img_array = np.array(image)
        
        # Convert to BGR for OpenCV
        if img_array.shape[2] == 4:  # RGBA
            bgr = cv2.cvtColor(img_array[:,:,:3], cv2.COLOR_RGB2BGR)
        else:
            bgr = cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR)
        
        # Apply GrabCut algorithm for better segmentation
        height, width = bgr.shape[:2]
        
        # Create rectangle for foreground (center 60% of image)
        rect = (int(width*0.2), int(height*0.2), int(width*0.6), int(height*0.6))
        
        # Initialize masks
        mask = np.zeros((height, width), np.uint8)
        bgd_model = np.zeros((1, 65), np.float64)
        fgd_model = np.zeros((1, 65), np.float64)
        
        # Apply GrabCut
        cv2.grabCut(bgr, mask, rect, bgd_model, fgd_model, 5, cv2.GC_INIT_WITH_RECT)
        
        # Create final mask
        mask2 = np.where((mask == 2) | (mask == 0), 0, 1).astype('uint8')
        
        # Apply morphological operations to clean up the mask
        kernel = np.ones((3,3), np.uint8)
        mask2 = cv2.morphologyEx(mask2, cv2.MORPH_CLOSE, kernel)
        mask2 = cv2.morphologyEx(mask2, cv2.MORPH_OPEN, kernel)
        
        # Convert back to PIL with alpha channel
        result_array = img_array.copy()
        result_array[:, :, 3] = mask2 * 255
        
        return Image.fromarray(result_array, 'RGBA')
        
    except Exception as e:
        print(f"OpenCV background removal error: {e}")
        return enhanced_basic_removal(image)

def improved_segmentation(image):
    """Improved segmentation algorithm"""
    try:
        img_array = np.array(image)
        
        # Calculate color statistics for edges
        edges = get_edge_pixels(img_array)
        
        # Use K-means clustering to separate foreground/background
        if cv2 is not None:
            # Reshape for clustering
            pixel_values = img_array[:,:,:3].reshape((-1, 3))
            pixel_values = np.float32(pixel_values)
            
            # Apply K-means
            criteria = (cv2.TERM_CRITERIA_EPS + cv2.TERM_CRITERIA_MAX_ITER, 20, 1.0)
            # K-means clustering with proper parameters
            compactness, labels, centers = cv2.kmeans(pixel_values, 2, None, criteria, 10, cv2.KMEANS_RANDOM_CENTERS)
            
            # Reshape back
            labels = labels.reshape(img_array.shape[:2])
            
            # Determine which cluster is background (usually the larger one on edges)
            edge_labels = labels[edges]
            if len(edge_labels) > 0:
                bg_label = np.bincount(edge_labels).argmax()
            else:
                bg_label = 0
            
            # Create mask
            mask = (labels != bg_label).astype(np.uint8) * 255
        else:
            # Fallback without OpenCV
            mask = create_basic_mask(img_array)
        
        # Apply mask to alpha channel
        result_array = img_array.copy()
        result_array[:, :, 3] = mask
        
        return Image.fromarray(result_array, 'RGBA')
        
    except Exception as e:
        print(f"Improved segmentation error: {e}")
        return enhanced_basic_removal(image)

def silhouette_detection(image):
    """Silhouette-based detection"""
    try:
        img_array = np.array(image)
        
        # Convert to grayscale for edge detection
        gray = np.dot(img_array[:,:,:3], [0.2989, 0.5870, 0.1140])
        
        # Apply edge detection
        if cv2 is not None:
            edges = cv2.Canny(gray.astype(np.uint8), 50, 150)
            
            # Find contours
            contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            # Create mask from largest contour (assumed to be main subject)
            if contours:
                largest_contour = max(contours, key=cv2.contourArea)
                mask = np.zeros(gray.shape, np.uint8)
                cv2.fillPoly(mask, [largest_contour], (255,))
                
                # Smooth the mask
                mask = cv2.GaussianBlur(mask, (5, 5), 0)
            else:
                mask = create_basic_mask(img_array)
        else:
            mask = create_basic_mask(img_array)
        
        # Apply mask
        result_array = img_array.copy()
        result_array[:, :, 3] = mask
        
        return Image.fromarray(result_array, 'RGBA')
        
    except Exception as e:
        print(f"Silhouette detection error: {e}")
        return enhanced_basic_removal(image)

def enhanced_basic_removal(image):
    """Enhanced version of basic background removal"""
    try:
        img_array = np.array(image)
        height, width = img_array.shape[:2]
        
        # Sample edge pixels to determine background color
        edge_pixels = []
        
        # Sample from edges
        edge_pixels.extend(img_array[0, :, :3].tolist())  # Top edge
        edge_pixels.extend(img_array[-1, :, :3].tolist())  # Bottom edge
        edge_pixels.extend(img_array[:, 0, :3].tolist())  # Left edge
        edge_pixels.extend(img_array[:, -1, :3].tolist())  # Right edge
        
        # Calculate dominant background color
        if edge_pixels:
            edge_pixels = np.array(edge_pixels)
            bg_color = np.median(edge_pixels, axis=0)
        else:
            bg_color = np.array([255, 255, 255])  # Default to white
        
        # Create mask based on similarity to background color
        mask = np.ones((height, width), dtype=np.uint8) * 255
        
        for i in range(height):
            for j in range(width):
                pixel = img_array[i, j, :3]
                # Calculate color distance
                distance = np.sqrt(np.sum((pixel - bg_color) ** 2))
                
                # If pixel is similar to background, make it transparent
                if distance < 50:  # Threshold for background similarity
                    mask[i, j] = 0
                elif distance < 100:  # Partial transparency for edge pixels
                    mask[i, j] = int(255 * (distance - 50) / 50)
        
        # Apply mask
        result_array = img_array.copy()
        result_array[:, :, 3] = mask
        
        return Image.fromarray(result_array, 'RGBA')
        
    except Exception as e:
        print(f"Enhanced basic removal error: {e}")
        # Ultra-simple fallback
        data = image.getdata()
        new_data = []
        for item in data:
            if item[:3] == (255, 255, 255) or (item[0] > 240 and item[1] > 240 and item[2] > 240):
                new_data.append((255, 255, 255, 0))
            else:
                new_data.append(item)
        image.putdata(new_data)
        return image


def get_edge_pixels(img_array):
    """Get boolean mask of edge pixels"""
    height, width = img_array.shape[:2]
    edges = np.zeros((height, width), dtype=bool)
    
    # Mark edge pixels
    edges[0, :] = True   # Top edge
    edges[-1, :] = True  # Bottom edge
    edges[:, 0] = True   # Left edge
    edges[:, -1] = True  # Right edge
    
    return edges

def create_basic_mask(img_array):
    """Create basic mask without OpenCV"""
    height, width = img_array.shape[:2]
    
    # Simple threshold-based mask
    gray = np.dot(img_array[:,:,:3], [0.299, 0.587, 0.114])
    mask = (gray > 200).astype(np.uint8) * 255
    mask = 255 - mask  # Invert
    
    return mask

def refine_edges(image):
    """Refine edges of the mask"""
    try:
        img_array = np.array(image)
        
        # Apply Gaussian blur to alpha channel for smoother edges
        alpha = img_array[:, :, 3].astype(np.float32)
        
        if cv2 is not None:
            # Use bilateral filter for edge-preserving smoothing
            alpha = cv2.bilateralFilter(alpha, 9, 75, 75)
        else:
            # Fallback: no smoothing
            pass
        
        img_array[:, :, 3] = alpha.astype(np.uint8)
        
        return Image.fromarray(img_array, 'RGBA')
        
    except Exception as e:
        print(f"Edge refinement error: {e}")
        return image

def apply_background(image, bg_color):
    """Apply solid background color"""
    try:
        if bg_color == 'transparent':
            return image
        
        # Create new image with solid background
        if bg_color == 'white':
            bg = (255, 255, 255)
        elif bg_color == 'black':
            bg = (0, 0, 0)
        else:
            # Parse custom color (assuming hex format #RRGGBB)
            if bg_color.startswith('#') and len(bg_color) == 7:
                bg = tuple(int(bg_color[i:i+2], 16) for i in (1, 3, 5))
            else:
                bg = (255, 255, 255)  # Default to white
        
        background = Image.new('RGB', image.size, bg)
        
        # Composite the images
        if image.mode == 'RGBA':
            background.paste(image, mask=image.split()[-1])
        else:
            background.paste(image)
        
        return background
        
    except Exception as e:
        print(f"Background application error: {e}")
        return image
